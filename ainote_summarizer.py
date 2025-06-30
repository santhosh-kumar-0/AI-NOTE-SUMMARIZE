import google.generativeai as genai
import fitz # For PDF
import pptx # For PPTX
from PIL import Image # For Images
import docx # For DOCX
import openpyxl # For XLSX
import pandas as pd # For CSV and better XLSX handling
from striprtf.striprtf import rtf_to_text # For RTF

import speech_recognition as sr # For Voice Input
import os # For temporary file management (less critical now, but still good to have)

from PyQt5.QtWidgets import (QWidget, QVBoxLayout, QTextEdit, QPushButton,
                             QLabel, QMessageBox, QFileDialog, QHBoxLayout, QInputDialog, 
                             QApplication, QLineEdit, QSizePolicy, QStackedLayout)
from PyQt5.QtGui import QFont, QPixmap, QImage # Added QImage for Pillow conversion
from PyQt5.QtCore import Qt, QTimer, QBuffer, QIODevice # Added QBuffer, QIODevice for Pillow conversion

from ui_styles import AppStyles
from voice_recognizer import VoiceRecognizer

# Constants for better readability and maintainability
MAX_INPUT_LENGTH = 1000000
MAX_OUTPUT_TOKENS = 8192

class AINoteSummarizer(QWidget):
    def __init__(self, stacked_widget): # Added stacked_widget parameter
        super().__init__()
        self.stacked_widget = stacked_widget # Store stacked_widget for logout
        self.setWindowTitle("✨ Keypoint AI 💡 - Main")
        self.setGeometry(90, 90, 900, 750)

        self.is_dark_theme = False # Will be set by LoginPage
        self.api_key = None
        self.voice_recognizer = VoiceRecognizer()
        self.current_image = None # Stores PIL Image object for summarization

        try:
            genai.configure(api_key="")
        except Exception:
            pass

        # --- Background Image Layer ---
        self.background_label = QLabel(self)
        self.background_label.setScaledContents(True) # Scale image to fit label size

        # --- Content Layer (Existing UI elements) ---
        content_v_layout = QVBoxLayout() # Renamed from main_layout for clarity in this stacked context
        content_v_layout.setSpacing(20)
        content_v_layout.setContentsMargins(30, 30, 30, 30)

        # Header with Title and Theme Toggle
        header_layout = QHBoxLayout()
        title_label = QLabel("✨ Keypoint AI 💡")
        title_label.setObjectName("SummarizerTitle")
        title_label.setFont(QFont("Segoe UI", 24, QFont.Bold))
        title_label.setAlignment(Qt.AlignLeft | Qt.AlignVCenter)

        self.api_settings_button = QPushButton("⚙️ API Settings", self)
        self.api_settings_button.setFixedSize(150, 40)
        self.api_settings_button.clicked.connect(self.show_api_key_dialog)
        self.api_settings_button.setFont(QFont("Segoe UI", 10, QFont.Bold))

        self.theme_toggle_button = QPushButton("🌙 Dark Mode", self)
        self.theme_toggle_button.setFixedSize(150, 40)
        self.theme_toggle_button.clicked.connect(self.toggle_theme) # Connected here
        self.theme_toggle_button.setFont(QFont("Segoe UI", 10, QFont.Bold))

        # Logout Button
        self.logout_button = QPushButton("Logout", self)
        self.logout_button.setFixedSize(80, 40)
        self.logout_button.clicked.connect(self.logout)
        self.logout_button.setFont(QFont("Segoe UI", 10, QFont.Bold))


        header_layout.addWidget(title_label)
        header_layout.addStretch(1)
        header_layout.addWidget(self.api_settings_button)
        header_layout.addWidget(self.theme_toggle_button)
        header_layout.addWidget(self.logout_button) # Add logout button to header
        content_v_layout.addLayout(header_layout) # Add to content_v_layout

        # --- Your Notes Section ---
        input_section_label = QLabel("Your Notes:")
        input_section_label.setObjectName("YourNotesLabel")
        input_section_label.setFont(QFont("Segoe UI", 16, QFont.Bold))
        input_section_label.setAlignment(Qt.AlignLeft)
        content_v_layout.addWidget(input_section_label)

        # --- New: Image Display Label ---
        self.image_display_label = QLabel(self) # Re-added image display label setup
        self.image_display_label.setAlignment(Qt.AlignCenter) # Re-added image display label setup
        self.image_display_label.setMinimumHeight(150) # Re-added image display label setup
        self.image_display_label.setMaximumHeight(300) # Re-added image display label setup
        self.image_display_label.setScaledContents(True) # Re-added image display label setup
        self.image_display_label.hide() # Initially hidden
        content_v_layout.addWidget(self.image_display_label) # Re-added image display label setup


        self.note_input = QTextEdit(self)
        self.note_input.setPlaceholderText("Type or paste your notes here, or upload a document...")
        self.note_input.setMinimumHeight(180)
        self.note_input.setFont(QFont("Segoe UI", 11))
        self.note_input.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        content_v_layout.addWidget(self.note_input)

        # --- Middle Buttons Layout: Voice Input, Upload, Summarize
        middle_buttons_layout = QHBoxLayout()
        middle_buttons_layout.setSpacing(15)

        self.voice_input_button = QPushButton("🎙️ Start Voice Input", self) # Renamed from record_button for consistency
        self.voice_input_button.setMinimumHeight(45)
        self.voice_input_button.setFont(QFont("Segoe UI", 12, QFont.Bold))
        self.voice_input_button.clicked.connect(self.start_voice_input)
        self.voice_input_button.setSizePolicy(QSizePolicy.MinimumExpanding, QSizePolicy.Preferred)
        middle_buttons_layout.addWidget(self.voice_input_button)

        self.upload_button = QPushButton("📁 Upload File", self)
        self.upload_button.setMinimumHeight(45)
        self.upload_button.setFont(QFont("Segoe UI", 12, QFont.Bold))
        self.upload_button.clicked.connect(self.upload_document)
        self.upload_button.setSizePolicy(QSizePolicy.MinimumExpanding, QSizePolicy.Preferred)
        middle_buttons_layout.addWidget(self.upload_button)

        self.summarize_button = QPushButton("✨ Summarize ", self)
        self.summarize_button.setMinimumHeight(55)
        self.summarize_button.setFont(QFont("Segoe UI", 14, QFont.Bold))
        self.summarize_button.clicked.connect(self.summarize_content)
        self.summarize_button.setSizePolicy(QSizePolicy.MinimumExpanding, QSizePolicy.Preferred)
        middle_buttons_layout.addWidget(self.summarize_button)

        content_v_layout.addLayout(middle_buttons_layout)


        # --- Summary Section ---
        output_section_label = QLabel("Summary:")
        output_section_label.setObjectName("SummaryLabel")
        output_section_label.setFont(QFont("Segoe UI", 16, QFont.Bold))
        output_section_label.setAlignment(Qt.AlignLeft)
        content_v_layout.addWidget(output_section_label)

        self.summary_output = QTextEdit(self)
        self.summary_output.setReadOnly(True)
        self.summary_output.setMinimumHeight(180)
        self.summary_output.setFont(QFont("Segoe UI", 11))
        self.summary_output.setLineWrapMode(QTextEdit.WidgetWidth)
        self.summary_output.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        content_v_layout.addWidget(self.summary_output)

        # --- Bottom Buttons Layout: Export and Clear All
        bottom_action_buttons_layout = QHBoxLayout()
        bottom_action_buttons_layout.setSpacing(15)

        # Changed to export as TXT
        self.export_button = QPushButton("💾 Export Summary as TXT", self)
        self.export_button.setMinimumHeight(45)
        self.export_button.setFont(QFont("Segoe UI", 12, QFont.Bold))
        self.export_button.clicked.connect(self.export_to_txt) # Connect to new export_to_txt method
        self.export_button.setSizePolicy(QSizePolicy.MinimumExpanding, QSizePolicy.Preferred)
        bottom_action_buttons_layout.addWidget(self.export_button)

        self.clear_button = QPushButton("🗑️ Clear All", self)
        self.clear_button.setMinimumHeight(45)
        self.clear_button.setFont(QFont("Segoe UI", 12, QFont.Bold))
        self.clear_button.clicked.connect(self.clear_all_inputs)
        self.clear_button.setSizePolicy(QSizePolicy.MinimumExpanding, QSizePolicy.Preferred)
        bottom_action_buttons_layout.addWidget(self.clear_button)

        content_v_layout.addLayout(bottom_action_buttons_layout)

        # --- Wrap content_v_layout in a QWidget to add to QStackedLayout ---
        content_widget = QWidget()
        content_widget.setLayout(content_v_layout)

        # --- Stacked Layout to layer background and content ---
        main_stacked_layout = QStackedLayout()
        main_stacked_layout.addWidget(self.background_label) # Layer 0: Background image
        main_stacked_layout.addWidget(content_widget)       # Layer 1: Actual content
        main_stacked_layout.setStackingMode(QStackedLayout.StackAll) # Show all layers

        self.setLayout(main_stacked_layout) # Set the stacked layout for the AINoteSummarizer page
        self.apply_theme_styles()

    def _set_background_image(self, image_path):
        """Loads and sets the background image to the QLabel."""
        try:
            pixmap = QPixmap(image_path)
            if pixmap.isNull():
                raise FileNotFoundError(f"Failed to load image: {image_path}. Pixmap is null.")
            self.background_label.setPixmap(pixmap)
        except Exception as e:
            # Fallback to a solid background color if image loading fails
            self.background_label.setStyleSheet("background-color: #1a1a2e;") # Dark fallback color
            print(f"Error loading main background image: {e}")

    def apply_theme_styles(self):
        self._set_background_image("images\main_background.jpg") # Set the main app background image here

        # Get the content_widget from the stacked layout to apply its background style
        content_widget = self.layout().itemAt(1).widget()

        if self.is_dark_theme:
            content_widget.setStyleSheet("background-color: rgba(26, 26, 46, 0.9); border-radius: 10px;") # Dark transparent
        else:
            content_widget.setStyleSheet("background-color: rgba(248, 249, 250, 0.9); border-radius: 10px;") # Light transparent

        # Apply other styles
        self.findChild(QLabel, "SummarizerTitle").setStyleSheet(AppStyles.get_title_style(self.is_dark_theme))
        self.findChild(QLabel, "YourNotesLabel").setStyleSheet(AppStyles.get_label_style(self.is_dark_theme, "16px", "bold"))
        self.findChild(QLabel, "SummaryLabel").setStyleSheet(AppStyles.get_label_style(self.is_dark_theme, "16px", "bold"))

        self.note_input.setStyleSheet(AppStyles.get_input_style(self.is_dark_theme))
        self.summary_output.setStyleSheet(AppStyles.get_input_style(self.is_dark_theme))

        self.upload_button.setStyleSheet(AppStyles.get_secondary_button_style(self.is_dark_theme))
        self.clear_button.setStyleSheet(AppStyles.get_secondary_button_style(self.is_dark_theme))
        self.voice_input_button.setStyleSheet(AppStyles.get_secondary_button_style(self.is_dark_theme))
        self.api_settings_button.setStyleSheet(AppStyles.get_secondary_button_style(self.is_dark_theme))
        self.export_button.setStyleSheet(AppStyles.get_secondary_button_style(self.is_dark_theme))
        self.theme_toggle_button.setStyleSheet(AppStyles.get_toggle_button_style(self.is_dark_theme))
        self.logout_button.setStyleSheet(AppStyles.get_secondary_button_style(self.is_dark_theme))

        self.summarize_button.setStyleSheet(AppStyles.get_primary_button_style(self.is_dark_theme))

        self.theme_toggle_button.setText("☀️ Light Mode" if self.is_dark_theme else "🌙 Dark Mode")

    def toggle_theme(self):
        """Toggles the current theme between dark and light and applies the new styles."""
        self.is_dark_theme = not self.is_dark_theme
        self.apply_theme_styles()

    def show_api_key_dialog(self):
        current_api_key = self.api_key or ""
        key, ok = QInputDialog.getText(self, 'API Key Settings', 'Enter your Gemini API key:', QLineEdit.Normal, current_api_key)
        if ok and key:
            self.api_key = key
            try:
                genai.configure(api_key=self.api_key)
                QMessageBox.information(self, "Success", "API key updated successfully!")
            except Exception as e:
                self.api_key = None
                QMessageBox.warning(self, "Error", f"Invalid API key format. Please ensure it's correct.\nError: {str(e)}")
        elif ok and not key:
            self.api_key = None
            QMessageBox.warning(self, "API Key Warning", "API key cannot be empty. Summarization might fail.")

    def clear_all_inputs(self):
        self.note_input.clear()
        self.current_image = None
        self.summary_output.clear()
        self.note_input.setPlaceholderText("Type or paste your notes here, or upload a document...")
        self.image_display_label.clear() # Clear the image from the label
        self.image_display_label.hide() # Hide the image label

    def convert_pil_to_qpixmap(self, pil_image): # Re-added this helper method
        """Converts a PIL Image object to a QPixmap."""
        if pil_image.mode == "RGB":
            # Convert to RGB if not already
            pil_image = pil_image.convert("RGB")
            # Convert PIL image to QImage
            qimage = QImage(pil_image.tobytes(), pil_image.size[0], pil_image.size[1],
                             QImage.Format_RGB888)
        elif pil_image.mode == "RGBA":
            # Convert PIL image to QImage (RGBA format)
            qimage = QImage(pil_image.tobytes(), pil_image.size[0], pil_image.size[1],
                             QImage.Format_RGBA8888)
        else:
            # For other modes, convert to RGBA first for broader compatibility
            pil_image = pil_image.convert("RGBA")
            qimage = QImage(pil_image.tobytes(), pil_image.size[0], pil_image.size[1],
                             QImage.Format_RGBA8888)

        return QPixmap.fromImage(qimage)

    def upload_document(self):
        # Define the file filter string directly
        file_filter = "All Supported Files (*.pdf *.pptx *.png *.jpg *.jpeg *.txt *.docx *.rtf *.xlsx *.csv);;PDF Files (*.pdf);;PowerPoint Files (*.pptx);;Image Files (*.png *.jpg *.jpeg);;Text Files (*.txt);;Word Documents (*.docx);;Rich Text Files (*.rtf);;Excel Files (*.xlsx);;CSV Files (*.csv)"
        
        # Call QFileDialog.getOpenFileName only ONCE
        file_path, selected_filter_name = QFileDialog.getOpenFileName(self, "Open File", "", file_filter)
        
        if file_path:
            self.current_image = None # Clear any previously loaded image for Gemini
            self.note_input.clear() # Clear note input when a new file is loaded
            self.image_display_label.clear() # Clear previous image from display
            self.image_display_label.hide() # Hide image display by default
            self.summary_output.setPlainText("Processing file...")
            QApplication.processEvents() # Update UI to show processing message

            file_extension = file_path.lower().split('.')[-1]

            if file_extension in ["png", "jpg", "jpeg"]:
                try:
                    self.current_image = Image.open(file_path) # Load image with Pillow for Gemini
                    qpixmap = self.convert_pil_to_qpixmap(self.current_image) # Convert for display
                    self.image_display_label.setPixmap(qpixmap) # Set image on the label
                    self.image_display_label.show() # Show the label

                    self.note_input.setPlainText(f"Image loaded: {file_path.split('/')[-1]}\n\n"
                                                 "Click 'Summarize' to get a visual summary.")
                    self.summary_output.setPlainText("Image loaded. Ready to summarize visual content.")
                except Exception as e:
                    QMessageBox.warning(self, "Image Load Error", f"Could not load image: {str(e)}")
                    self.current_image = None
                    self.image_display_label.clear()
                    self.image_display_label.hide()
                    self.summary_output.setPlainText("Failed to load image.")
            else:
                extracted_text = self.extract_text_from_file(file_path)
                self.note_input.setPlainText(extracted_text)
                self.summary_output.setPlainText("Text extraction complete. You can now summarize.")

    def extract_text_from_file(self, file_path):
        try:
            file_extension = file_path.lower().split('.')[-1]
            text = ""

            if file_extension == "pdf":
                try:
                    doc = fitz.open(file_path)
                    text = "\n".join([page.get_text("text") for page in doc])
                    doc.close()
                except ImportError:
                    QMessageBox.critical(self, "Missing Dependency",
                                         "PyMuPDF (fitz) is not installed. Please install it using: pip install PyMuPDF")
                    return ""
                except Exception as pdf_e:
                    QMessageBox.warning(self, "PDF Read Error",
                                         f"Could not read PDF file. It might be corrupted or encrypted: {pdf_e}")
                    return ""
            elif file_extension == "pptx":
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + "\n"
            elif file_extension == "txt":
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            elif file_extension == "docx":
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            elif file_extension == "rtf":
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    rtf_content = f.read()
                try:
                    text = rtf_to_text(rtf_content)
                except Exception as rtf_e:
                    QMessageBox.warning(self, "RTF Error", f"Could not parse RTF content. May contain unsupported elements.\nError: {rtf_e}")
                    text = rtf_content # Fallback to raw RTF if parsing fails
            elif file_extension == "xlsx":
                workbook = openpyxl.load_workbook(file_path)
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value is not None:
                                text += str(cell.value) + " "
                        text += "\n"
            elif file_extension == "csv":
                df = pd.read_csv(file_path)
                text = df.to_string(index=False)
            else:
                QMessageBox.warning(self, "Unsupported File Type", f"File type '.{file_extension}' is not supported for text extraction.")
                return ""

            return text.strip()
        except Exception as e:
            QMessageBox.warning(self, "Error Extracting Text", f"An unexpected error occurred while extracting text from {file_path.split('/')[-1]}: {str(e)}")
            return ""

    def start_voice_input(self):
        self.voice_input_button.setEnabled(False) # Renamed from record_button for consistency
        self.original_note_text_before_voice = self.note_input.toPlainText()
        self.note_input.setHtml("🎙️ <b>Listening... Please speak clearly.</b>")
        self.note_input.append("_Say something or wait for timeout if you are done._")
        self.current_image = None # Clear any loaded image if starting voice input
        self.image_display_label.clear() # Clear displayed image
        self.image_display_label.hide() # Hide image display
        QApplication.processEvents()

        QTimer.singleShot(100, self._process_voice_input)

    def _process_voice_input(self):
        recognized_text, error = self.voice_recognizer.listen_and_recognize()

        if recognized_text:
            if "🎙️ <b>Listening" in self.note_input.toHtml():
                self.note_input.setPlainText(self.original_note_text_before_voice + "\n\n" + recognized_text)
            else:
                self.note_input.append("\n\n" + recognized_text)
            QMessageBox.information(self, "Voice Input", "Voice input complete. Text added to your notes.")
        elif error:
            self.note_input.setPlainText(self.original_note_text_before_voice)
            QMessageBox.warning(self, "Voice Input Error",
                                 f"{error}\n\n**Tips:**\n"
                                 f"- Ensure your microphone is properly connected and selected.\n"
                                 f"- Speak clearly and at a moderate pace.\n"
                                 f"- Minimize background noise.\n"
                                 f"- Check your internet connection for online recognition."
                                )
        else:
            self.note_input.setPlainText(self.original_note_text_before_voice)
            QMessageBox.warning(self, "Voice Input Error", "No input detected or unrecognized speech. Please try again.")

        self.voice_input_button.setEnabled(True) # Renamed from record_button for consistency

    def summarize_content(self):
        if not self.api_key:
            self.summary_output.setPlainText("Gemini API Key is not set. Please go to '⚙️ API Settings' to configure it.")
            QMessageBox.warning(self, "API Key Missing", "Please set your Gemini API Key in the API Settings.")
            return

        genai.configure(api_key=self.api_key)
        model = genai.GenerativeModel("gemini-1.5-flash")

        self.summary_output.setPlainText("Generating summary with Gemini 1.5 Flash...")
        QApplication.processEvents()

        try:
            contents = []
            if self.current_image:
                prompt = "Please provide a concise summary and description of this image, identifying key objects, scenes, and any visible text. Aim for clarity and conciseness, and structure the summary in bullet points or short paragraphs."
                contents = [prompt, self.current_image] # Pass Pillow Image object directly
            else:
                note_text = self.note_input.toPlainText().strip()
                if not note_text:
                    self.summary_output.setPlainText("Please enter a note, upload a document, or load an image to summarize.")
                    return

                if len(note_text) > MAX_INPUT_LENGTH:
                    self.summary_output.setPlainText(
                        f"Input too long ({len(note_text)} characters).\n"
                        f"Maximum recommended length for optimal performance: {MAX_INPUT_LENGTH} characters.\n"
                        "Please shorten your text or split it into smaller parts."
                    )
                    return

                prompt = (f"Please provide a comprehensive summary of the following text, "
                          f"including key points and main ideas. Aim for clarity and conciseness, "
                          f"and and structure the summary with bullet points or short paragraphs:\n\n{note_text}")
                contents = [prompt]

            if contents:
                response = model.generate_content(
                    contents,
                    generation_config={
                        "max_output_tokens": MAX_OUTPUT_TOKENS,
                        "temperature": 0.3,
                    }
                )

                if response.text:
                    summary = response.text.strip()
                    self.summary_output.setPlainText(summary)
                else:
                    self.summary_output.setPlainText("No summary was generated. The AI might not have found enough content or encountered an internal issue.")

                self.summary_output.verticalScrollBar().setValue(0)
            else:
                self.summary_output.setPlainText("No content to summarize.")

        except Exception as e:
            error_message = str(e)
            if "quota" in error_message.lower():
                self.summary_output.setPlainText(f"Error: You have exceeded your API quota. Please wait or check your Google Cloud Console for details.\n\n{error_message}")
                QMessageBox.critical(self, "API Quota Exceeded", "You have exceeded your Gemini API quota. Please try again later or check your billing details on Google Cloud.")
            elif "authentication" in error_message.lower() or "api key" in error_message.lower():
                self.summary_output.setPlainText(f"Error: API Key authentication failed. Please verify your Gemini API Key in settings.\n\n{error_message}")
                QMessageBox.critical(self, "API Key Error", "Authentication failed. Please check your Gemini API Key.")
            else:
                self.summary_output.setPlainText(f"An unexpected error occurred during summarization: {error_message}\n"
                                                 "Please check your internet connection or try a shorter text.")
                QMessageBox.critical(self, "Summarization Error", f"An unexpected error occurred: {e}")
        finally:
            pass

    def export_to_txt(self): # Renamed from export_to_pdf
        summary_text = self.summary_output.toPlainText()
        if not summary_text or "Error:" in summary_text or "No summary" in summary_text or "Generating summary" in summary_text:
            QMessageBox.warning(self, "Export Failed", "No valid summary to export. Please generate a summary first.")
            return

        # Changed to save as TXT
        file_path, _ = QFileDialog.getSaveFileName(self, "Save Summary as Text", "", "Text Files (*.txt)")
        if file_path:
            try:
                # Ensure the file has a .txt extension
                if not file_path.lower().endswith('.txt'):
                    file_path += '.txt'
                    
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write(summary_text)
                QMessageBox.information(self, "Success", "Summary exported as TXT successfully!")
            except Exception as e:
                QMessageBox.warning(self, "Export Error", f"Failed to export summary as TXT.\nError: {str(e)}")

    def logout(self):
        """Logs out the user and returns to the login page."""
        self.stacked_widget.setCurrentIndex(0) # Go back to Login Page
        self.note_input.clear() # Clear input for next session
        self.summary_output.clear() # Clear output
        self.current_image = None # Clear any loaded image
        self.image_display_label.clear() # Clear the image from the label
        self.image_display_label.hide() # Hide the image label
        # Optionally, reset theme to default for login page:
        login_page = self.stacked_widget.widget(0)
        login_page.is_dark_theme = False # Ensure login page starts in light mode
        login_page.apply_theme_styles()